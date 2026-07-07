import os
from pptx import Presentation

path = r'E:\Github\260707_SG_AI를_이용한_피착재_분석_및_매칭과_제품_역설계_통합_서비스_개발.pptx'
prs = Presentation(path)

# --- Slide 12 (0-indexed slide index 11) ---
slide12 = prs.slides[11]
for shape in slide12.shapes:
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs)
        if "임시 방어" in text:
            # Replace paragraph 0 text
            p0 = shape.text_frame.paragraphs[0]
            p0.text = "임시 방어 로직에서 GPU 가속 인프라로 격상"
            # Add new paragraphs
            new_lines = [
                "  - Nvidia CUDA 가속 컨테이너 전환 (003, 007)",
                "  - docker-compose.yml 내 GPU 자원 할당(Reservation) 구성 완료",
                "  - 고해상도 실시간 3D 재구성 병목 현상 완전 해소"
            ]
            for line in new_lines:
                p = shape.text_frame.add_paragraph()
                p.text = line

# --- Slide 23 (0-indexed slide index 22) ---
slide23 = prs.slides[22]
for shape in slide23.shapes:
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs)
        if "SG_proj_004" in text:
            shape.text_frame.paragraphs[0].text = "SG_proj_004 (개발 및 테스트 격리 완료)"
        elif "PostgreSQL" in text or "SQLite" in text or "FastAPI" in text:
            shape.text_frame.add_paragraph()
            p = shape.text_frame.add_paragraph()
            p.text = "인메모리 테스트 격리: respx 및 sqlite:///:memory:를 연동하여 E2E 테스트 실행 시 운영 DB 오염 방지 안전망 구축"

# --- Slide 31 (0-indexed slide index 30) ---
slide31 = prs.slides[30]
for shape in slide31.shapes:
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs)
        if "SG_proj_011" in text:
            shape.text_frame.paragraphs[0].text = "SG_proj_011 (개발 및 패널티 연동 완료)"
        elif "가공성" in text or "경계값" in text or "비접촉" in text:
            shape.text_frame.add_paragraph()
            p = shape.text_frame.add_paragraph()
            p.text = "기재별 패널티 연동: configs/config.json에 기재 종류(PVC, PE, PO) 및 두께별 수축 패널티 테이블을 정의하고, 011 결과에 가감산 적용"

# --- Slide 32 (0-indexed slide index 31) ---
slide32 = prs.slides[31]
for shape in slide32.shapes:
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs)
        if "SG_proj_012" in text:
            shape.text_frame.paragraphs[0].text = "SG_proj_012 (개발 및 스펙 이원화 완료)"
        elif "물성" in text or "매칭" in text or "수동" in text:
            shape.text_frame.add_paragraph()
            p = shape.text_frame.add_paragraph()
            p.text = "점착력 이원화 대응: target_adhesion 스펙을 초기 점착력(target_initial_adhesion) 및 후기 점착력(target_aged_adhesion)으로 분리 연동"

# --- Slide 33 (0-indexed slide index 32) ---
slide33 = prs.slides[32]
for shape in slide33.shapes:
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs)
        if "SG_proj_013" in text:
            shape.text_frame.paragraphs[0].text = "SG_proj_013 (개발 및 FTIR 피드백 연동 완료)"
        elif "오차" in text or "xgboost" in text or "confidence" in text or "피드백" in text:
            shape.text_frame.add_paragraph()
            p = shape.text_frame.add_paragraph()
            p.text = "FTIR 스펙트럼 피드백: 009의 transmittance(투과율) 파형 데이터를 ir_gnn_features 입력 특징량으로 정합하여 013 보정 알고리즘 실체화"

# --- Slide 34 (0-indexed slide index 33) ---
slide34 = prs.slides[33]
for shape in slide34.shapes:
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs)
        if "SG_proj_014" in text:
            shape.text_frame.paragraphs[0].text = "SG_proj_014 (개발 및 도메인 검증 완료)"
        elif "중앙 제어" in text or "비동기" in text or "E2E" in text or "예외" in text:
            shape.text_frame.add_paragraph()
            p = shape.text_frame.add_paragraph()
            p.text = "도메인 불변 규칙 단속: schemas.py에서 Pydantic validator를 통한 초기 점착력 <= 후기 점착력 및 target_tg <= 0.0°C 규칙 물리적 강제화"

# --- Slide 38 (0-indexed slide index 37) ---
slide38 = prs.slides[37]
for shape in slide38.shapes:
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs)
        if "SG_integration_step3" in text:
            shape.text_frame.paragraphs[0].text = "SG_integration_step3 (통합 및 E2E 실증 완료)"
        elif "001(" in text or "역설계" in text or "분자" in text or "배합" in text:
            # Replace target text inside paragraphs
            for p in shape.text_frame.paragraphs:
                if "연계 구조 설계" in p.text or "최종 물성 조율" in p.text:
                    p.text = "역설계 및 고분자 조율: 자사(001), 오픈(006), IR(009-n=10 중합도 모사) 예측 모델과 최종 물성 조율(013)의 연계 인터페이스 구축 완료"
            shape.text_frame.add_paragraph()
            p = shape.text_frame.add_paragraph()
            p.text = "CI/CD 자동 검증: GitHub Actions CI 파이프라인(Ruff, mypy, pytest) 연동 완료로 E2E 안전망 가동"

prs.save(path)
print("Successfully modified PPTX slides with latest sprint outcomes.")
