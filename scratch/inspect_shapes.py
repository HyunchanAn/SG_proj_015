from pptx import Presentation

path = r'E:\Github\260707_SG_AI를_이용한_피착재_분석_및_매칭과_제품_역설계_통합_서비스_개발.pptx'
prs = Presentation(path)

slides_to_inspect = [12, 23, 31, 32, 33, 34, 38]

for slide_idx in slides_to_inspect:
    # 0-indexed slide index (Slide 12 is index 11)
    slide = prs.slides[slide_idx - 1]
    print(f"\n=== Slide {slide_idx} ===")
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"  Shape {shape_idx}:")
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                print(f"    P {p_idx}: '{p.text}'")
